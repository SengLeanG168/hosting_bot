import telebot
import pytesseract
from PIL import Image
import io
import os
# For PDF, Word, PowerPoint
import fitz  # PyMuPDF
from pdf2image import convert_from_bytes
import docx
from pptx import Presentation
from docx import Document
from docx2pdf import convert as docx2pdf_convert
import tempfile
import shutil

# ដាក់ Tesseract executable path (ប្រសិនបើមិនមាននៅក្នុង system PATH)
# ឧទាហរណ៍សម្រាប់ Windows:
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

# ជំនួស 'YOUR_BOT_TOKEN' ជាមួយ Token របស់ Bot អ្នក
API_TOKEN = '7328885744:AAGDvDt85Se9oenNL-tAxr9MIMkU7ytuN30'
bot = telebot.TeleBot(API_TOKEN)

# មុខងារសម្រាប់ command /start
@bot.message_handler(commands=['start'])
def send_welcome(message):
    welcome_message = """
សូមស្វាគមន៍មកកាន់ Telegram bot ដែលបម្លែងពីរូបភាពទៅជាអក្សរ ។​ \n
Welcome to the image to text conversion Telegram bot.\n
ខ្ញុំបាទឈ្មោះ ហេង សេងលៀង ជាអ្នកបង្កើត Telegram bot នេះឡើងក្នុងគោលបំណងសម្រាប់ងាយស្រួលដល់និស្សិតធ្វើការបម្លែងពីរូបភាពទៅជាអក្សរ ដូច្នេះសូមប្រើដោយក្ដីរីករាយ ! \n
I am Heng Sengleang, the creator of this Telegram bot, designed to facilitate students in converting images to text. Please use it with joy!\n
សូមផ្ញើរូបភាពមកខ្ញុំដើម្បីធ្វើការបម្លែងពីរូបភាពទៅជាអក្សរ ។​​​ \n
Please send an image to me for conversion from image to text.
"""
    bot.reply_to(message, welcome_message)

# មុខងារសម្រាប់ទទួលរូបភាព

# Handle images, PDFs, Word, PowerPoint
@bot.message_handler(content_types=['photo', 'document'])
def handle_file(message):
    try:
        processing_message = bot.reply_to(message, "កំពុងបម្លែងឯកសារទៅជាអក្សរ...")

        file_info = None
        file_bytes = None
        file_name = None
        mime_type = None

        if message.content_type == 'photo':
            file_info = bot.get_file(message.photo[-1].file_id)
            file_bytes = bot.download_file(file_info.file_path)
            file_name = 'image.jpg'
            mime_type = 'image/jpeg'
        elif message.content_type == 'document':
            file_info = bot.get_file(message.document.file_id)
            file_bytes = bot.download_file(file_info.file_path)
            file_name = message.document.file_name.lower()
            mime_type = message.document.mime_type
        else:
            bot.edit_message_text(chat_id=message.chat.id, message_id=processing_message.message_id, text="សូមផ្ញើរឯកសារដែលគាំទ្រ (រូបភាព, PDF, Word, PowerPoint)")
            return

        # Handle image
        if mime_type and mime_type.startswith('image'):
            img = Image.open(io.BytesIO(file_bytes))
            text = pytesseract.image_to_string(img, lang='khm+eng', config='--psm 6')
            if text.strip():
                bot.edit_message_text(chat_id=message.chat.id, message_id=processing_message.message_id, text=text)
            else:
                bot.edit_message_text(chat_id=message.chat.id, message_id=processing_message.message_id, text="មិនអាចស្គាល់អក្សរបានទេ សូមផ្ញើរូបភាពដែលច្បាស់ជាងមុន។")
            return


        # PDF to Word (docx)
        if file_name.endswith('.pdf') and (message.caption and 'toword' in message.caption.lower()):
            try:
                with tempfile.TemporaryDirectory() as tmpdir:
                    pdf_path = os.path.join(tmpdir, 'input.pdf')
                    docx_path = os.path.join(tmpdir, 'output.docx')
                    # Write PDF and ensure file is closed before opening with fitz
                    with open(pdf_path, 'wb') as f:
                        f.write(file_bytes)
                    # Now open with fitz (PyMuPDF)
                    pdf_doc = fitz.open(pdf_path)
                    doc = Document()
                    for page in pdf_doc:
                        text = page.get_text()
                        doc.add_paragraph(text)
                    pdf_doc.close()
                    doc.save(docx_path)
                    with open(docx_path, 'rb') as f:
                        bot.send_document(message.chat.id, f, visible_file_name='converted.docx')
                bot.edit_message_text(chat_id=message.chat.id, message_id=processing_message.message_id, text="បម្លែង PDF ទៅ Word ជោគជ័យ។")
            except Exception as e:
                bot.edit_message_text(chat_id=message.chat.id, message_id=processing_message.message_id, text=f"បរាជ័យ: {e}")
            return

        # Word (docx) to PDF
        if file_name.endswith('.docx') and (message.caption and 'topdf' in message.caption.lower()):
            try:
                with tempfile.TemporaryDirectory() as tmpdir:
                    docx_path = os.path.join(tmpdir, 'input.docx')
                    pdf_path = os.path.join(tmpdir, 'output.pdf')
                    with open(docx_path, 'wb') as f:
                        f.write(file_bytes)
                    docx2pdf_convert(docx_path, pdf_path)
                    with open(pdf_path, 'rb') as f:
                        bot.send_document(message.chat.id, f, visible_file_name='converted.pdf')
                bot.edit_message_text(chat_id=message.chat.id, message_id=processing_message.message_id, text="បម្លែង Word ទៅ PDF ជោគជ័យ។")
            except Exception as e:
                bot.edit_message_text(chat_id=message.chat.id, message_id=processing_message.message_id, text=f"បរាជ័យ: {e}")
            return

        # PowerPoint (pptx) to PowerPoint (copy)
        if file_name.endswith('.pptx') and (message.caption and 'topptx' in message.caption.lower()):
            try:
                with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp:
                    tmp.write(file_bytes)
                    tmp.flush()
                    with open(tmp.name, 'rb') as f:
                        bot.send_document(message.chat.id, f, visible_file_name='copy.pptx')
                bot.edit_message_text(chat_id=message.chat.id, message_id=processing_message.message_id, text="បម្លែង PowerPoint ជោគជ័យ (ចម្លង)។")
            except Exception as e:
                bot.edit_message_text(chat_id=message.chat.id, message_id=processing_message.message_id, text=f"បរាជ័យ: {e}")
            return

        # PDF to PDF (copy)
        if file_name.endswith('.pdf') and (message.caption and 'topdf' in message.caption.lower()):
            try:
                with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp:
                    tmp.write(file_bytes)
                    tmp.flush()
                    with open(tmp.name, 'rb') as f:
                        bot.send_document(message.chat.id, f, visible_file_name='copy.pdf')
                bot.edit_message_text(chat_id=message.chat.id, message_id=processing_message.message_id, text="ចម្លង PDF ជោគជ័យ។")
            except Exception as e:
                bot.edit_message_text(chat_id=message.chat.id, message_id=processing_message.message_id, text=f"បរាជ័យ: {e}")
            return

        # Handle Word (docx) to text (if not converting)
        if file_name.endswith('.docx') or (mime_type and 'word' in mime_type):
            try:
                doc = docx.Document(io.BytesIO(file_bytes))
                text = '\n'.join([para.text for para in doc.paragraphs])
                if text.strip():
                    bot.edit_message_text(chat_id=message.chat.id, message_id=processing_message.message_id, text=text[:4096])
                else:
                    bot.edit_message_text(chat_id=message.chat.id, message_id=processing_message.message_id, text="មិនមានអត្ថបទក្នុងឯកសារ Word។")
            except Exception as e:
                bot.edit_message_text(chat_id=message.chat.id, message_id=processing_message.message_id, text=f"កំហុស Word: {e}")
            return

        # Handle PowerPoint (pptx)
        if file_name.endswith('.pptx') or (mime_type and 'presentation' in mime_type):
            try:
                prs = Presentation(io.BytesIO(file_bytes))
                text = ''
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + '\n'
                if text.strip():
                    bot.edit_message_text(chat_id=message.chat.id, message_id=processing_message.message_id, text=text[:4096])
                else:
                    bot.edit_message_text(chat_id=message.chat.id, message_id=processing_message.message_id, text="មិនមានអត្ថបទក្នុង PowerPoint។")
            except Exception as e:
                bot.edit_message_text(chat_id=message.chat.id, message_id=processing_message.message_id, text=f"កំហុស PowerPoint: {e}")
            return

        # Unsupported file
        bot.edit_message_text(chat_id=message.chat.id, message_id=processing_message.message_id, text="សូមផ្ញើរឯកសារដែលគាំទ្រ (រូបភាព, PDF, Word, PowerPoint)")
    except Exception as e:
        error_message = f"មានកំហុស៖ {e}"
        bot.edit_message_text(chat_id=message.chat.id, message_id=processing_message.message_id, text=error_message)
        print(f"Error occurred: {e}")

# ចាប់ផ្ដើម Bot
print("Bot is running...")
bot.polling(none_stop=True)